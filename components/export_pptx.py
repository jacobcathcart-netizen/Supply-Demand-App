"""Export scenario results to a branded CCR PowerPoint presentation.

Generates a multi-slide deck with KPI metrics, chart images, and
enterprise styling using the CCR brand palette.
"""

from __future__ import annotations

from datetime import date
from io import BytesIO
from pathlib import Path
from typing import Any

from matplotlib.figure import Figure
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from components.branding import (
    GRAY,
    LIGHT_BLUE,
    LIGHT_GRAY,
    NAVY,
    ORANGE,
    WARM_WHITE,
    WHITE,
    YELLOW,
)

# ── Constants ──────────────────────────────────────────────────────────

_SLIDE_WIDTH = Inches(13.333)
_SLIDE_HEIGHT = Inches(7.5)

_LOGO_PATH = Path(__file__).resolve().parent.parent / "assets" / "logo.jpg"

_NAVY = RGBColor.from_string(NAVY.lstrip("#"))
_LIGHT_BLUE = RGBColor.from_string(LIGHT_BLUE.lstrip("#"))
_ORANGE = RGBColor.from_string(ORANGE.lstrip("#"))
_YELLOW = RGBColor.from_string(YELLOW.lstrip("#"))
_WHITE = RGBColor.from_string("FFFFFF")
_GRAY = RGBColor.from_string(GRAY.lstrip("#"))
_LIGHT_GRAY = RGBColor.from_string(LIGHT_GRAY.lstrip("#"))
_WARM_WHITE = RGBColor.from_string(WARM_WHITE.lstrip("#").replace("#", ""))

_FONT_FAMILY = "Tahoma"


# ── Helpers ────────────────────────────────────────────────────────────


def _hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert a hex color string (#RRGGBB) to an RGBColor."""
    return RGBColor.from_string(hex_color.lstrip("#"))


def _set_slide_bg(slide, color: RGBColor) -> None:
    """Set a solid background color on a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_shape(
    slide,
    left: int,
    top: int,
    width: int,
    height: int,
    fill_color: RGBColor | None = None,
    border_color: RGBColor | None = None,
    border_width: Pt | None = None,
) -> Any:
    """Add a rectangle shape to a slide."""
    from pptx.enum.shapes import MSO_SHAPE

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False

    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()

    if border_color and border_width:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()

    return shape


def _add_text_box(
    slide,
    left: int,
    top: int,
    width: int,
    height: int,
    text: str,
    font_size: int = 12,
    font_color: RGBColor = _NAVY,
    bold: bool = False,
    alignment: PP_ALIGN = PP_ALIGN.LEFT,
    font_name: str = _FONT_FAMILY,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
) -> Any:
    """Add a text box with styled text to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    txBox.text_frame.auto_size = None
    txBox.text_frame.paragraphs[0].alignment = alignment

    # Vertical alignment
    txBox.text_frame.paragraphs[0].space_before = Pt(0)
    txBox.text_frame.paragraphs[0].space_after = Pt(0)

    run = txBox.text_frame.paragraphs[0].runs[0] if txBox.text_frame.paragraphs[0].runs else txBox.text_frame.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold
    run.font.name = font_name

    return txBox


def _fig_to_image_stream(fig: Figure, dpi: int = 200) -> BytesIO:
    """Render a matplotlib Figure to a PNG BytesIO stream."""
    buf = BytesIO()
    fig.savefig(buf, format="png", dpi=dpi, bbox_inches="tight", facecolor=fig.get_facecolor())
    buf.seek(0)
    return buf


def _add_logo(slide, left: int, top: int, height: int = Inches(0.6)) -> None:
    """Add the CCR logo to a slide if available."""
    if _LOGO_PATH.exists():
        slide.shapes.add_picture(str(_LOGO_PATH), left, top, height=height)


def _add_footer_bar(slide, text: str = "") -> None:
    """Add a thin branded footer bar at the bottom of the slide."""
    bar_height = Inches(0.35)
    bar_top = _SLIDE_HEIGHT - bar_height
    _add_shape(slide, 0, bar_top, _SLIDE_WIDTH, bar_height, fill_color=_NAVY)

    if text:
        _add_text_box(
            slide,
            Inches(0.5),
            bar_top + Inches(0.04),
            _SLIDE_WIDTH - Inches(1),
            bar_height - Inches(0.08),
            text,
            font_size=8,
            font_color=_WHITE,
            alignment=PP_ALIGN.RIGHT,
        )


# ── Slide builders ─────────────────────────────────────────────────────


def _build_title_slide(
    prs: Presentation,
    scenario_name: str,
    region_label: str,
    report_date: str,
) -> None:
    """Slide 1: Full-bleed title slide with CCR branding."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    _set_slide_bg(slide, _NAVY)

    # Accent bar at top
    _add_shape(slide, 0, 0, _SLIDE_WIDTH, Inches(0.08), fill_color=_LIGHT_BLUE)

    # Logo
    _add_logo(slide, Inches(0.7), Inches(0.6), height=Inches(1.0))

    # Title
    _add_text_box(
        slide,
        Inches(0.7),
        Inches(2.2),
        Inches(10),
        Inches(1.2),
        "Supply & Demand",
        font_size=44,
        font_color=_WHITE,
        bold=True,
    )
    _add_text_box(
        slide,
        Inches(0.7),
        Inches(3.2),
        Inches(10),
        Inches(0.8),
        "Scenario Results",
        font_size=36,
        font_color=_LIGHT_BLUE,
        bold=True,
    )

    # Divider line
    _add_shape(
        slide,
        Inches(0.7),
        Inches(4.2),
        Inches(3),
        Inches(0.04),
        fill_color=_YELLOW,
    )

    # Subtitle info
    _add_text_box(
        slide,
        Inches(0.7),
        Inches(4.6),
        Inches(8),
        Inches(0.4),
        scenario_name,
        font_size=18,
        font_color=_WHITE,
        bold=False,
    )
    _add_text_box(
        slide,
        Inches(0.7),
        Inches(5.15),
        Inches(8),
        Inches(0.4),
        f"Region: {region_label}",
        font_size=14,
        font_color=_hex_to_rgb("#8EAED0"),
    )
    _add_text_box(
        slide,
        Inches(0.7),
        Inches(5.6),
        Inches(8),
        Inches(0.4),
        report_date,
        font_size=14,
        font_color=_hex_to_rgb("#8EAED0"),
    )

    # Bottom accent
    _add_shape(
        slide,
        0,
        _SLIDE_HEIGHT - Inches(0.08),
        _SLIDE_WIDTH,
        Inches(0.08),
        fill_color=_YELLOW,
    )


def _build_kpi_slide(
    prs: Presentation,
    metrics: dict[str, dict[str, str]],
) -> None:
    """Slide 2: KPI dashboard with metric cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, _hex_to_rgb(WARM_WHITE))

    # Header bar
    _add_shape(slide, 0, 0, _SLIDE_WIDTH, Inches(0.9), fill_color=_NAVY)
    _add_text_box(
        slide,
        Inches(0.6),
        Inches(0.15),
        Inches(8),
        Inches(0.6),
        "Key Performance Indicators",
        font_size=26,
        font_color=_WHITE,
        bold=True,
    )
    _add_logo(slide, _SLIDE_WIDTH - Inches(1.8), Inches(0.12), height=Inches(0.65))

    # Section: Supply & Demand
    _add_text_box(
        slide,
        Inches(0.6),
        Inches(1.2),
        Inches(4),
        Inches(0.4),
        "SUPPLY & DEMAND",
        font_size=12,
        font_color=_LIGHT_BLUE,
        bold=True,
    )

    supply_demand_keys = [
        "Baseline Supply",
        "Scenario Supply",
        "Total Demand",
        "Supply Delta",
    ]
    _draw_metric_cards(slide, supply_demand_keys, metrics, top=Inches(1.7))

    # Section: Gap & Backlog
    _add_text_box(
        slide,
        Inches(0.6),
        Inches(4.0),
        Inches(4),
        Inches(0.4),
        "GAP & BACKLOG",
        font_size=12,
        font_color=_LIGHT_BLUE,
        bold=True,
    )

    gap_backlog_keys = [
        "Baseline Gap",
        "Scenario Gap",
        "Initial Backlog",
        "Ending Backlog",
    ]
    _draw_metric_cards(slide, gap_backlog_keys, metrics, top=Inches(4.5))

    _add_footer_bar(slide, "Cypress Creek Renewables | Workforce Planning")


def _draw_metric_cards(
    slide,
    keys: list[str],
    metrics: dict[str, dict[str, str]],
    top: int,
) -> None:
    """Draw a row of 4 metric cards."""
    card_width = Inches(2.8)
    card_height = Inches(1.8)
    margin_left = Inches(0.6)
    gap = Inches(0.25)

    for i, key in enumerate(keys):
        m = metrics.get(key, {"value": "N/A", "delta": ""})
        left = margin_left + i * (card_width + gap)

        # Card background
        card = _add_shape(
            slide, left, top, card_width, card_height, fill_color=_WHITE
        )
        # Left accent border
        _add_shape(
            slide, left, top, Inches(0.06), card_height, fill_color=_LIGHT_BLUE
        )

        # Label
        _add_text_box(
            slide,
            left + Inches(0.25),
            top + Inches(0.2),
            card_width - Inches(0.4),
            Inches(0.3),
            key.upper(),
            font_size=9,
            font_color=_GRAY,
            bold=True,
        )

        # Value
        _add_text_box(
            slide,
            left + Inches(0.25),
            top + Inches(0.6),
            card_width - Inches(0.4),
            Inches(0.6),
            m["value"],
            font_size=22,
            font_color=_NAVY,
            bold=True,
        )

        # Delta (if present)
        if m.get("delta"):
            delta_color = _hex_to_rgb("#007647") if "+" not in m["delta"] else _ORANGE
            # For supply delta, + is good; for backlog delta, - is good
            _add_text_box(
                slide,
                left + Inches(0.25),
                top + Inches(1.25),
                card_width - Inches(0.4),
                Inches(0.35),
                m["delta"],
                font_size=10,
                font_color=_GRAY,
                bold=False,
            )


def _build_chart_slide(
    prs: Presentation,
    title: str,
    fig: Figure,
    caption: str = "",
) -> None:
    """Add a slide with a full-width chart image."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, _hex_to_rgb(WARM_WHITE))

    # Header bar
    _add_shape(slide, 0, 0, _SLIDE_WIDTH, Inches(0.9), fill_color=_NAVY)
    _add_text_box(
        slide,
        Inches(0.6),
        Inches(0.15),
        Inches(10),
        Inches(0.6),
        title,
        font_size=26,
        font_color=_WHITE,
        bold=True,
    )
    _add_logo(slide, _SLIDE_WIDTH - Inches(1.8), Inches(0.12), height=Inches(0.65))

    # Chart image — fill the available area
    img_stream = _fig_to_image_stream(fig)
    chart_top = Inches(1.1)
    chart_left = Inches(0.3)
    chart_width = _SLIDE_WIDTH - Inches(0.6)
    chart_height = Inches(5.4)

    slide.shapes.add_picture(img_stream, chart_left, chart_top, chart_width, chart_height)

    # Caption
    if caption:
        _add_text_box(
            slide,
            Inches(0.6),
            Inches(6.6),
            _SLIDE_WIDTH - Inches(1.2),
            Inches(0.4),
            caption,
            font_size=9,
            font_color=_GRAY,
        )

    _add_footer_bar(slide, "Cypress Creek Renewables | Workforce Planning")


def _build_dual_chart_slide(
    prs: Presentation,
    title: str,
    fig_top: Figure,
    fig_bottom: Figure,
    caption: str = "",
) -> None:
    """Add a slide with two stacked chart images (used for sensitivity)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, _hex_to_rgb(WARM_WHITE))

    # Header bar
    _add_shape(slide, 0, 0, _SLIDE_WIDTH, Inches(0.9), fill_color=_NAVY)
    _add_text_box(
        slide,
        Inches(0.6),
        Inches(0.15),
        Inches(10),
        Inches(0.6),
        title,
        font_size=26,
        font_color=_WHITE,
        bold=True,
    )
    _add_logo(slide, _SLIDE_WIDTH - Inches(1.8), Inches(0.12), height=Inches(0.65))

    chart_left = Inches(0.3)
    chart_width = _SLIDE_WIDTH - Inches(0.6)

    # Top chart
    img_top = _fig_to_image_stream(fig_top)
    slide.shapes.add_picture(
        img_top, chart_left, Inches(1.05), chart_width, Inches(2.9)
    )

    # Bottom chart
    img_bottom = _fig_to_image_stream(fig_bottom)
    slide.shapes.add_picture(
        img_bottom, chart_left, Inches(4.05), chart_width, Inches(2.9)
    )

    if caption:
        _add_text_box(
            slide,
            Inches(0.6),
            Inches(6.95),
            _SLIDE_WIDTH - Inches(1.2),
            Inches(0.3),
            caption,
            font_size=9,
            font_color=_GRAY,
        )

    _add_footer_bar(slide, "Cypress Creek Renewables | Workforce Planning")


# ── Public API ─────────────────────────────────────────────────────────


def build_presentation(
    *,
    scenario_name: str,
    region_label: str,
    metrics: dict[str, dict[str, str]],
    fig_baseline: Figure | None = None,
    fig_scenario: Figure | None = None,
    fig_gap: Figure | None = None,
    fig_backlog: Figure | None = None,
    fig_sensitivity_fan: Figure | None = None,
    fig_sensitivity_tornado: Figure | None = None,
) -> bytes:
    """Build a complete branded PowerPoint and return raw bytes.

    Parameters
    ----------
    scenario_name:
        Display name for the scenario.
    region_label:
        Region filter label (e.g. "All Selected Regions" or "NC - Central").
    metrics:
        Dict mapping metric names to ``{"value": "...", "delta": "..."}``
        dicts.  Expected keys: Baseline Supply, Scenario Supply,
        Total Demand, Supply Delta, Baseline Gap, Scenario Gap,
        Initial Backlog, Ending Backlog.
    fig_*:
        Matplotlib Figure objects for each chart tab.  ``None`` charts
        are skipped.

    Returns
    -------
    bytes
        The ``.pptx`` file content, ready for ``st.download_button``.
    """
    prs = Presentation()

    # Set widescreen slide dimensions (13.333 x 7.5 inches)
    prs.slide_width = _SLIDE_WIDTH
    prs.slide_height = _SLIDE_HEIGHT

    today = date.today().strftime("%B %d, %Y")

    # 1. Title slide
    _build_title_slide(prs, scenario_name, region_label, today)

    # 2. KPI dashboard
    _build_kpi_slide(prs, metrics)

    # 3. Chart slides
    if fig_baseline is not None:
        _build_chart_slide(
            prs,
            "Baseline Supply vs Demand",
            fig_baseline,
            caption="Baseline supply and demand over time with gap analysis.",
        )

    if fig_scenario is not None:
        _build_chart_slide(
            prs,
            "Scenario Supply vs Demand",
            fig_scenario,
            caption="Scenario supply and demand with headcount adjustments applied.",
        )

    if fig_gap is not None:
        _build_chart_slide(
            prs,
            "Gap Analysis",
            fig_gap,
            caption="Monthly gap comparison: baseline vs scenario.",
        )

    if fig_backlog is not None:
        _build_chart_slide(
            prs,
            "Backlog Trend",
            fig_backlog,
            caption="Cumulative backlog trend with normalized backlog (squad-months).",
        )

    # 4. Sensitivity — dual chart or single
    if fig_sensitivity_fan is not None and fig_sensitivity_tornado is not None:
        _build_dual_chart_slide(
            prs,
            "Sensitivity Analysis",
            fig_sensitivity_fan,
            fig_sensitivity_tornado,
            caption="Top: backlog sensitivity envelope.  Bottom: parameter impact tornado.",
        )
    elif fig_sensitivity_fan is not None:
        _build_chart_slide(
            prs,
            "Sensitivity Analysis",
            fig_sensitivity_fan,
            caption="Backlog sensitivity envelope showing the range of possible outcomes.",
        )
    elif fig_sensitivity_tornado is not None:
        _build_chart_slide(
            prs,
            "Sensitivity Analysis",
            fig_sensitivity_tornado,
            caption="Tornado chart showing each parameter's impact on ending backlog.",
        )

    # Serialize to bytes
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()
